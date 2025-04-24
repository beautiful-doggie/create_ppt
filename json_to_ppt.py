import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import unicodedata
import json 
# 配置常量
MAX_CHARS_PER_SLIDE = 300  # 每页最大字数
BACKGROUND_COLOR = RGBColor(34, 45, 65)  # 深蓝背景
TITLE_COLOR = RGBColor(255, 255, 255)  # 白色标题
TEXT_COLOR = RGBColor(255, 255, 255)  # 白色文本
FONT_NAME = 'Microsoft YaHei'  # 字体

# 清理文本函数
def clean_text(text):
    """清除空行 + 清理特殊字符 + 连接句子"""
    def remove_invisible(s):
        s = s.replace('\r', '')  # 移除回车
        s = s.replace('\x0b', '')  # 垂直制表符
        s = s.replace('\x0c', '')  # 换页符
        s = s.replace('\u200b', '')  # 零宽空格
        return ''.join(c for c in s if unicodedata.category(c)[0] != 'C')  # 去掉控制字符

    lines = [remove_invisible(line.strip()) for line in text.strip().split('\n') if line.strip()]
    return '，'.join(lines)

# 估算文本长度
def estimate_weighted_length(text):
    """估算文本长度，考虑换行符对排版的影响"""
    base_len = len(text.replace('\n', ''))
    newline_count = text.count('\n')
    return base_len + newline_count * 10  # 每个换行大约占一行高度

# 合并小段文本
def merge_blocks(blocks, max_chars=MAX_CHARS_PER_SLIDE):
    """合并小段，避免换行撑爆页面"""
    merged = []
    buffer = ''
    for block in blocks:
        block = block.strip()
        test_block = buffer + '\n' + block if buffer else block
        if estimate_weighted_length(test_block) <= max_chars:
            buffer = test_block
        else:
            if buffer:
                merged.append(buffer.strip())
            buffer = block

    if buffer:
        merged.append(buffer.strip())
    return merged

# 智能分段（应用部分）
def split_application(text):
    """按常见编号样式智能分段 + 合并为合理分页"""
    pattern = r'(?=\n?(\d+[\.、]|[一二三四五六七八九十]+、|（[一二三四五六七八九十]+）|[①②③④⑤⑥⑦⑧⑨⑩]))'
    cleaned = text.strip()
    raw_blocks = re.split(pattern, cleaned)

    # 组合分段（re.split 保留了编号在奇数位上）
    combined_blocks = []
    i = 0
    while i < len(raw_blocks):
        if i + 1 < len(raw_blocks):
            combined_blocks.append(raw_blocks[i] + raw_blocks[i+1])
            i += 2
        else:
            combined_blocks.append(raw_blocks[i])
            i += 1

    # 清理每个分段，确保没有空内容
    combined_blocks = [b.strip() for b in combined_blocks if b.strip() and len(b.strip()) > 1]

    # 如果有一些分段依旧是空的（例如仅包含编号），我们跳过
    return merge_blocks(combined_blocks)


# 智能分段（概念部分）
def split_concept(text):
    """句子合并为自然段"""
    lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
    # 每两三行为一个段落
    buffer = ''
    raw_blocks = []
    for line in lines:
        buffer += line + '，'
        if len(buffer) > 60:
            raw_blocks.append(buffer.strip('，') + '。')
            buffer = ''
    if buffer:
        raw_blocks.append(buffer.strip('，') + '。')
    return merge_blocks(raw_blocks)

# 智能分段（比较部分）
def split_comparison(text):
    """按 vs 或换行分 + 合并"""
    items = re.split(r'\n| vs | VS ', text)
    raw_blocks = []
    buffer = ''
    for item in items:
        item = item.strip()
        if not item:
            continue
        buffer += item + '，'
        if len(buffer) > 60:
            raw_blocks.append(buffer.strip('，') + '。')
            buffer = ''
    if buffer:
        raw_blocks.append(buffer.strip('，') + '。')
    return merge_blocks(raw_blocks)

# 智能分段（代码部分）
def split_code(text):
    """按注释划分代码段 + 合并小段"""
    lines = text.strip().split('\n')
    raw_blocks = []
    current_block = []

    for line in lines:
        if re.match(r'^\s*(#|//)', line):  # 是注释开头
            if current_block:
                raw_blocks.append('\n'.join(current_block))
            current_block = [line]
        else:
            current_block.append(line)

    if current_block:
        raw_blocks.append('\n'.join(current_block))

    return merge_blocks(raw_blocks)

# 按字段类型分段
def split_field(text, field_type):
    if field_type == "application":
        return split_application(text)
    elif field_type == "concept":
        return split_concept(text)
    elif field_type == "comparison":
        return split_comparison(text)
    elif field_type == "code":
        return split_code(text)
    else:
        return merge_blocks([clean_text(text)])

# 设置字体样式
def set_font_style(paragraph, font_size=Pt(18), font_name=FONT_NAME, font_color=TEXT_COLOR):
    run = paragraph.add_run()
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    return paragraph
import os 
# 添加PPT页面
def add_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容布局

    # 设置标题
    title_box = slide.shapes.title
    title_box.text = title
    title_frame = title_box.text_frame.paragraphs[0]
    title_frame.font.size = Pt(24)
    title_frame.font.color.rgb = TITLE_COLOR
    title_frame.font.name = FONT_NAME

    # 设置背景色
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    # 设置内容框
    content_box = slide.shapes.placeholders[1]
    content_frame = content_box.text_frame
    content_frame.clear()

    # 添加段落文本（不处理图片）
    paragraphs = content.strip().split('\n')
    for i, para_text in enumerate(paragraphs):
        para = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        para.text = para_text.strip()
        para.font.color.rgb = TEXT_COLOR
        para.font.size = Pt(18)
        set_font_style(para, font_size=Pt(18))

# 添加目录页
def add_table_of_contents(prs, content_keys):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录 Contents"

    content_box = slide.shapes.placeholders[1]
    frame = content_box.text_frame
    frame.clear()

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    for key in content_keys:
        para = frame.add_paragraph()
        para.text = f"{key.capitalize()}"
        para.font.size = Pt(20)
        para.font.color.rgb = TEXT_COLOR
        set_font_style(para, font_size=Pt(20))

# 添加章节页
def add_section_slide(prs, key):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"📘 {key.capitalize()} Section"
    slide.placeholders[1].text = ""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR


def add_illustration_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    para = title_frame.paragraphs[0]
    para.font.size = Pt(24)
    para.font.color.rgb = TITLE_COLOR
    para.font.name = FONT_NAME

    # 插入图片（如存在）
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.5), width=Inches(7))
import re
from collections import Counter

def find_best_matching_image(image_dir, segments):
    # 合并所有段落内容
    full_text = " ".join(segments)

    # 简单方式：按词语分词（这里你可以用jieba分词更准确）
    words = re.findall(r'[\u4e00-\u9fa5a-zA-Z0-9]+', full_text)  # 中英文词语
    word_freq = Counter(words)

    # 获取所有图片名
    image_files = [f for f in os.listdir(image_dir) if f.lower().endswith(".png")]

    # 统计图片名命中频率
    image_scores = []
    for img_file in image_files:
        basename = os.path.splitext(img_file)[0]
        score = sum(word_freq.get(word, 0) for word in words if word in basename)
        image_scores.append((score, img_file))

    image_scores.sort(reverse=True)  # 按匹配度降序排列
    best_score, best_image = image_scores[0] if image_scores else (0, None)

    if best_score > 0:
        return os.path.join(image_dir, best_image)
    else:
        return None
# 从JSON生成PPT
def json_to_ppt(topic):
    json_path = f"jsons/{topic}.json"
    ppt_path = f"ppt/{topic}.pptx"
    
    if not os.path.exists(json_path):
        print(f"未找到 {json_path}")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    topic = data.get("topic", "未知主题")
    content = data.get("content", {})

    prs = Presentation()

    # 封面页
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = "由 AI 自动生成"

    # 目录页
    content_keys = list(content.keys())
    add_table_of_contents(prs, content_keys)

    # 内容页
    for key in content_keys:
        text = content[key]
        segments = split_field(text, key)

        # 插图页（仅插入一次）
        image_path = find_best_matching_image("illustrations", segments)
        if image_path:
            illustration_title = f"{key.capitalize()} illustration"
            add_illustration_slide(prs, illustration_title, image_path)
        else:
            print(f"未找到匹配插图：字段 {key}")

        # 添加文字内容页
        for i, segment in enumerate(segments):
            slide_title = f"{key.capitalize()}（{i+1}）"
            add_slide(prs, slide_title, segment)

    prs.save(ppt_path)
    print(f"PPT 已生成：{ppt_path}")
if __name__=='__main__':
    json_to_ppt('栈')
# def json_to_ppt(json_path, output_path, image_dir="illustrations"):
#     with open(json_path, 'r', encoding='utf-8') as f:
#         data = json.load(f)

#     topic = data.get("topic", "未知主题")
#     content = data.get("content", {})

#     prs = Presentation()

#     # 封面页
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])
#     title_slide.shapes.title.text = topic
#     title_slide.placeholders[1].text = "由 AI 自动生成"

#     # 目录页
#     content_keys = list(content.keys())
#     add_table_of_contents(prs, content_keys)

#     # 内容页
#     for key in content_keys:
#         text = content[key]
#         segments = split_field(text, key)

#         # 插图页（仅插入一次）
#         image_path = find_best_matching_image(image_dir, segments)
#         if image_path:
#             illustration_title = f"{key.capitalize()} illustration"
#             add_illustration_slide(prs, illustration_title, image_path)
#         else:
#             print(f"未找到匹配插图：字段 {key}")

#         # 添加文字内容页
#         for i, segment in enumerate(segments):
#             slide_title = f"{key.capitalize()}（{i+1}）"
#             add_slide(prs, slide_title, segment)

#     prs.save(output_path)
#     print(f"PPT 已生成：{output_path}")

# # 示例运行
# if __name__ == "__main__":
#     json_to_ppt("jsons/二叉树.json", "output.pptx")
